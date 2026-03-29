"""PDF转换器 - PDF与Word/PPT/图片之间的转换"""
import os
from utils.logger import logger, log_conversion


def pdf_to_word(input_path: str, output_path: str) -> str:
    """PDF 转 Word (.docx)，使用 pdf2docx"""
    try:
        from pdf2docx import Converter

        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()

        log_conversion(input_path, output_path, True)
        return output_path
    except Exception as e:
        error_msg = str(e)
        log_conversion(input_path, output_path, False, error_msg)
        # 检测是否为扫描件
        if _is_likely_scanned(input_path):
            raise RuntimeError('检测到扫描PDF，转换效果可能较差') from e
        raise RuntimeError('转换失败，请检查文件是否损坏') from e


def pdf_to_ppt(input_path: str, output_path: str) -> str:
    """PDF 转 PPT (.pptx) - 使用 PyMuPDF 渲染每页为图片再插入PPT，无需 Poppler"""
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches
        import tempfile

        doc = fitz.open(input_path)
        prs = Presentation()
        # 设置为宽屏16:9
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # 空白布局
        mat = fitz.Matrix(200 / 72, 200 / 72)  # 200 dpi

        with tempfile.TemporaryDirectory() as tmp_dir:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_path = os.path.join(tmp_dir, f'page_{i}.png')
                pix.save(img_path)

                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    img_path,
                    Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height
                )

            prs.save(output_path)
        doc.close()

        log_conversion(input_path, output_path, True)
        return output_path
    except Exception as e:
        error_msg = str(e)
        log_conversion(input_path, output_path, False, error_msg)
        raise RuntimeError('PDF转PPT失败，请检查文件是否损坏') from e


def pdf_to_images(input_path: str, output_dir: str, fmt: str = 'png', dpi: int = 200) -> list:
    """PDF 转图片，每页一张；多页时在 output_dir 下创建以PDF命名的子目录。使用 PyMuPDF，无需 Poppler"""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(input_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        mat = fitz.Matrix(dpi / 72, dpi / 72)

        # 多页创建子目录，单页直接放在 output_dir
        if len(doc) > 1:
            save_dir = os.path.join(output_dir, base_name)
            os.makedirs(save_dir, exist_ok=True)
        else:
            save_dir = output_dir

        output_paths = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            out_name = f'{base_name}_page_{i + 1}.{fmt}'
            out_path = os.path.join(save_dir, out_name)
            if fmt.lower() == 'png':
                pix.save(out_path)
            else:
                # jpg / jpeg：通过 PIL 控制压缩质量
                from PIL import Image
                import io
                img = Image.open(io.BytesIO(pix.tobytes('png')))
                img.save(out_path, 'JPEG', quality=92)
            output_paths.append(out_path)

        doc.close()
        log_conversion(input_path, save_dir, True)
        return output_paths
    except Exception as e:
        error_msg = str(e)
        log_conversion(input_path, output_dir, False, error_msg)
        raise RuntimeError('PDF转图片失败，请检查文件是否损坏') from e


def _is_likely_scanned(pdf_path: str) -> bool:
    """简单判断PDF是否为扫描件（文本内容极少）"""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        total_text = ''
        for page in doc:
            total_text += page.get_text()
        doc.close()
        # 如果文本很少，可能是扫描件
        return len(total_text.strip()) < 50
    except Exception:
        return False
